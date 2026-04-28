from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, glob

# --- CONFIG ---
INDIGO = RGBColor(79, 70, 229)
PURPLE = RGBColor(139, 92, 246)
DARK = RGBColor(30, 30, 46)
GRAY = RGBColor(100, 100, 120)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(248, 247, 255)
ACCENT_GREEN = RGBColor(16, 185, 129)
W, H = Inches(13.333), Inches(7.5)  # 16:9

IMG_DIR = r"C:\Users\venky\.gemini\antigravity\brain\2c4ccb17-0134-44bd-9b3b-a16703f81abf"

def find_img(keyword):
    matches = glob.glob(os.path.join(IMG_DIR, f"{keyword}*.png"))
    return matches[0] if matches else None

def set_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=28, color=DARK, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_slide(prs, title, bullets, img_key=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)
    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), H)
    bar.fill.solid(); bar.fill.fore_color.rgb = INDIGO; bar.line.fill.background()
    # title
    add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.9), title, size=36, color=INDIGO, bold=True)
    # underline
    uline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.15), Inches(2), Inches(0.06))
    uline.fill.solid(); uline.fill.fore_color.rgb = PURPLE; uline.line.fill.background()

    if img_key:
        img_path = find_img(img_key)
        # bullets on left, image on right
        bx = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(6.5), Inches(5.5))
        tf = bx.text_frame; tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"▸  {b}"
            p.font.size = Pt(20); p.font.color.rgb = DARK; p.font.name = "Calibri"
            p.space_after = Pt(14)
        if img_path:
            slide.shapes.add_picture(img_path, Inches(7.5), Inches(1.5), Inches(5.3), Inches(5.3))
    else:
        bx = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(11.5), Inches(5.5))
        tf = bx.text_frame; tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"▸  {b}"
            p.font.size = Pt(22); p.font.color.rgb = DARK; p.font.name = "Calibri"
            p.space_after = Pt(16)

def main():
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H

    # ===== SLIDE 1: TITLE =====
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s1, RGBColor(240, 237, 255))
    img = find_img("slide1_title")
    if img: s1.shapes.add_picture(img, Inches(3.5), Inches(0.3), Inches(6), Inches(3.5))
    add_text(s1, Inches(0.5), Inches(4.0), Inches(12), Inches(1), "CreatorLink", size=52, color=INDIGO, bold=True, align=PP_ALIGN.CENTER)
    add_text(s1, Inches(0.5), Inches(5.0), Inches(12), Inches(0.6), "A Role-Driven Creator Marketplace Platform", size=26, color=PURPLE, align=PP_ALIGN.CENTER)
    add_text(s1, Inches(0.5), Inches(5.8), Inches(12), Inches(0.5), "D. Venka Reddy  (R210363)  •  Guide: Mr. Satyanandaram N  •  RGUKT RK Valley, Dept of CSE", size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # ===== SLIDE 2: PROBLEM =====
    add_bullet_slide(prs, "The Problem in the Creator Economy", [
        "Brands rely on informal channels (WhatsApp, email) to hire editors, photographers & influencers.",
        "Generalist platforms like Fiverr/Upwork don't enforce creative-industry specific rules.",
        "No built-in mechanism to transition from gig negotiation to a private project workspace.",
        "Independent creators lack a centralized portfolio and gig-specific rating system.",
        "Result: Fragmented workflows, miscommunication, and wasted time for everyone involved."
    ], "slide2_problem")

    # ===== SLIDE 3: SOLUTION =====
    add_bullet_slide(prs, "Introducing CreatorLink", [
        "A specialized full-stack marketplace tailored exclusively for the creative ecosystem.",
        "Supports 5 distinct user roles: Admin, Brand, Influencer, Editor, Photographer.",
        "Bridges the gap between Brands sourcing talent and Creators offering services.",
        "Strict Role-Based Access Control (RBAC) enforced at both API and UI layers.",
        "End-to-end workflow: Gig Posting → Competitive Bidding → Real-time Chat → Reviews."
    ], "slide3_roles")

    # ===== SLIDE 4: OBJECTIVES =====
    add_bullet_slide(prs, "Core Objectives", [
        "Build a Role-Based Access Control marketplace supporting 5 user types.",
        "Implement a structured gig posting and competitive bidding system.",
        "Automatically provision real-time WebSocket chat upon bid acceptance.",
        "Develop a robust review and 5-star rating system post-completion.",
        "Deliver a secure, high-performance architecture using modern web technologies.",
        "Ensure seamless user experience with lazy loading, dark mode, and responsive design."
    ])

    # ===== SLIDE 5: WORKFLOW =====
    add_bullet_slide(prs, "Platform Workflow", [
        "1. DISCOVER  —  Brand posts a Gig. Creators filter & search the marketplace.",
        "2. NEGOTIATE  —  Creators submit proposals with amount, turnaround & message.",
        "3. ACCEPT  —  Brand accepts a bid. All competing bids are auto-rejected.",
        "4. COLLABORATE  —  Real-time WebSocket chat workspace is automatically created.",
        "5. COMPLETE & REVIEW  —  Work is delivered, marked done, both parties leave ratings."
    ], "slide4_workflow")

    # ===== SLIDE 6: RBAC =====
    add_bullet_slide(prs, "Role-Based Access Control (RBAC)", [
        "ADMIN  —  Platform moderation: manage users, gigs, bids. Cannot post or bid.",
        "BRAND  —  Posts gigs for Editing, Photography, or Promotion. Cannot bid.",
        "INFLUENCER  —  Bids on Promotion gigs from Brands. Can also post hiring gigs.",
        "EDITOR  —  Bids exclusively on Editing gigs. Cannot post gigs.",
        "PHOTOGRAPHER  —  Bids exclusively on Photography gigs. Cannot post gigs.",
        "Enforced at API level (FastAPI dependencies) AND UI level (conditional rendering)."
    ])

    # ===== SLIDE 7: TECH STACK =====
    add_bullet_slide(prs, "Technology Stack", [
        "Frontend:  React 19 + Vite + Tailwind CSS v4 + React Router v7 + Lucide Icons.",
        "Backend:  FastAPI (Python) + Uvicorn ASGI Server + Python-JOSE for JWT.",
        "Database:  MongoDB Atlas (Cloud) + Beanie ODM + Motor Async Driver.",
        "Media:  Cloudinary SDK for image/file uploads and CDN delivery.",
        "Real-time:  Native FastAPI WebSockets with ConnectionManager singleton.",
        "Auth:  JWT Access + Refresh tokens with automatic silent refresh on expiry."
    ], "slide5_architecture")

    # ===== SLIDE 8: KEY FEATURES =====
    add_bullet_slide(prs, "Key Features Implemented", [
        "Gig Marketplace with category filters, budget range, and sorting options.",
        "Competitive Bidding with amount, turnaround days, and proposal messages.",
        "Real-time WebSocket Chat with typing indicators and file sharing.",
        "Public Profiles with portfolio links, images, skills, and star ratings.",
        "Admin Dashboard with user/gig/bid moderation and role-based filtering.",
        "Lazy Loading (React.lazy + Suspense) for optimized initial load performance.",
        "Dark Cinematic UI theme with glassmorphism and micro-animations."
    ])

    # ===== SLIDE 9: DATABASE =====
    add_bullet_slide(prs, "Database Design", [
        "MongoDB Atlas — Perfect for nested data: portfolios, skills, chat participants.",
        "6 Core Collections:  Users, Gigs, Bids, Conversations, Messages, Reviews.",
        "Beanie ODM with Pydantic validation ensures strict schemas on NoSQL.",
        "Indexed queries on gig categories, user roles, and conversation participants.",
        "Async Motor driver for non-blocking database I/O operations."
    ])

    # ===== SLIDE 10: TESTING =====
    add_bullet_slide(prs, "Testing & Validation", [
        "Unit Testing:  100% pass rate on business logic, role checks, JWT handling.",
        "Integration Testing:  Full workflow (Post → Bid → Accept → Chat → Review).",
        "Security:  Tested Token Replay, JWT Tampering, NoSQL injection, unauthorized access.",
        "UAT:  5 test users, 100% task success rate, avg gig posting time: 3.5 minutes.",
        "API Testing:  All endpoints verified via Postman with edge case coverage."
    ])

    # ===== SLIDE 11: CHALLENGES =====
    add_bullet_slide(prs, "Engineering Challenges & Solutions", [
        "Challenge: Concurrent WebSocket connections crashing on disconnect.",
        "Solution: Built singleton ConnectionManager with JWT validation on handshake.",
        "Challenge: Complex role-based validation across 5 user types.",
        "Solution: Granular permission checks inside specific API route handlers.",
        "Challenge: 403 errors on page refresh due to token expiry race conditions.",
        "Solution: Axios interceptor with queued token refresh and silent retry."
    ])

    # ===== SLIDE 12: FUTURE SCOPE =====
    add_bullet_slide(prs, "Future Scope", [
        "Integrated Payments — Escrow model using Razorpay or Stripe.",
        "AI-Powered Matching — NLP to auto-recommend creators for gig descriptions.",
        "Video Portfolios — HLS adaptive bitrate streaming for video editors.",
        "Mobile App — React Native client extending the existing FastAPI backend.",
        "Analytics Dashboard — Revenue tracking, engagement metrics, conversion rates."
    ])

    # ===== SLIDE 13: CONCLUSION =====
    add_bullet_slide(prs, "Conclusion", [
        "Successfully delivered a modern, decoupled, three-tier web application.",
        "Solved a real-world problem for the creator economy with domain-specific tools.",
        "Strict RBAC ensures platform integrity across all 5 user roles.",
        "Real-time collaboration via WebSockets bridges the gap from hiring to delivery.",
        "Built a scalable, production-ready foundation for future enhancements."
    ])

    # ===== SLIDE 14: THANK YOU =====
    s14 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s14, RGBColor(240, 237, 255))
    img = find_img("slide6_closing")
    if img: s14.shapes.add_picture(img, Inches(3), Inches(0.5), Inches(7), Inches(4))
    add_text(s14, Inches(0.5), Inches(4.8), Inches(12), Inches(1), "Thank You!", size=52, color=INDIGO, bold=True, align=PP_ALIGN.CENTER)
    add_text(s14, Inches(0.5), Inches(5.8), Inches(12), Inches(0.6), "Questions & Answers", size=28, color=PURPLE, align=PP_ALIGN.CENTER)
    add_text(s14, Inches(0.5), Inches(6.5), Inches(12), Inches(0.5), "D. Venka Reddy  •  R210363  •  RGUKT RK Valley", size=16, color=GRAY, align=PP_ALIGN.CENTER)

    out = r"C:\Users\venky\OneDrive\Desktop\Creator Link\CreatorLink_Final_Presentation.pptx"
    prs.save(out)
    print(f"Presentation saved to: {out}")

if __name__ == "__main__":
    main()
